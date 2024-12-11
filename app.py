
from flask import Flask, render_template, request, send_file, jsonify,abort,redirect,url_for
import os
import shutil
from werkzeug.utils import secure_filename
import subprocess
import traceback
import tempfile
from PIL import Image,ImageDraw
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
from concurrent.futures import ThreadPoolExecutor, as_completed
import io
from language import init_babel, get_translation
from datetime import datetime
from TransDocs.Backend.functionality.recyclebin import recycle_bin_bp
#update 
#update 1

app = Flask(__name__)
app.register_blueprint(recycle_bin_bp)
init_babel(app)
app.config['UPLOAD_FOLDER'] = os.path.abspath('uploads')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024 * 1024  # 16GB max file size
# storing file details
file_data_store= {}
# Ensure upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


# Set the Poppler path (you may need to adjust this based on your system setup)
POPPLER_PATH = r"C:\poppler-24.07.0\Library\bin"
MAX_WORKERS = os.cpu_count()  # Use all available CPU cores
CHUNK_SIZE = 10  # Increased chunk size for better parallelism

def convert_to_png(input_file, output_file):
    try:
        subprocess.run(['pandoc', '-o', output_file, input_file], check=True)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Error converting {input_file} to PNG: {e}")
        return False
    except Exception as e:
        print(f"Unexpected error in convert_to_png: {e}")
        return False

def perform_ocr(image):
    try:
        return pytesseract.image_to_string(image)
    except Exception as e:
        print(f"Error performing OCR: {e}")
        return ""

def process_image_chunk(chunk):
    text = ""
    for i, image in enumerate(chunk):
        text += f"--- Page/Slide {i + 1} ---\n\n{perform_ocr(image)}\n\n"
        # Release memory
        image.close()
    return text

def process_images(images):
    chunks = [images[i:i + CHUNK_SIZE] for i in range(0, len(images), CHUNK_SIZE)]
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        results = list(executor.map(process_image_chunk, chunks))
    return "".join(results)

def process_pdf(input_file):
    try:
        with tempfile.TemporaryDirectory() as path:
            images = convert_from_path(input_file, output_folder=path, poppler_path=POPPLER_PATH)
            return process_images(images)
    except Exception as e:
        print(f"Error processing PDF: {e}")
        return ""

def process_pptx(input_file):
    try:
        prs = Presentation(input_file)
        images = []
        for slide in prs.slides:
            img = Image.new('RGB', (prs.slide_width, prs.slide_height), color='white')
            draw = ImageDraw.Draw(img)
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    draw.text((shape.left, shape.top), shape.text, fill='black')
            images.append(img)
        return process_images(images)
    except Exception as e:
        print(f"Error processing PPTX: {e}")
        return ""

def extract_text_from_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading text file {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading .docx file {file_path}: {e}")
        return ""

def extract_text_from_excel(file_path):
    try:
        with pd.ExcelFile(file_path) as xls:
            dfs = {sheet_name: pd.read_excel(xls, sheet_name) for sheet_name in xls.sheet_names}
        return "\n".join([f"\n--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}\n" for sheet_name, df in dfs.items()])
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {e}")
        return ""

def process_document(input_file):
    try:
        _, ext = os.path.splitext(input_file)
        ext = ext.lower()
        
        if ext == '.pdf':
            return process_pdf(input_file)
        elif ext == '.pptx':
            return process_pptx(input_file)
        elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            with Image.open(input_file) as img:
                return perform_ocr(img)
        elif ext in ['.txt', '.csv', '.json', '.xml', '.yaml', '.yml', '.ini', '.toml', '.py', '.js', '.html', '.htm', '.css']:
            return extract_text_from_file(input_file)
        elif ext == '.docx':
            return extract_text_from_docx(input_file)
        elif ext in ['.xls', '.xlsx', '.ods']:
            return extract_text_from_excel(input_file)
        else:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                temp_png = temp_file.name
            if convert_to_png(input_file, temp_png):
                with Image.open(temp_png) as img:
                    return perform_ocr(img)
            else:
                return "Error: Unable to process the file"
    except Exception as e:
        print(f"Error processing document: {e}")
        print(traceback.format_exc())
        return "Error: Unable to process the document"

@app.route('/')
def index():
      return render_template('index.html')

@app.route('/home')
def home():
    path = request.args.get('path', '')
    full_path = os.path.abspath(os.path.join(app.config['UPLOAD_FOLDER'], path))
    
    # Ensure the path is within the UPLOAD_FOLDER
    if not full_path.startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        abort(403, description="Access denied")
    
    if not os.path.exists(full_path) or not os.path.isdir(full_path):
        abort(404, description="Folder not found")
    
    files = []
    folders = []
    for item in os.listdir(full_path):
        item_path = os.path.join(full_path, item)
        if os.path.isfile(item_path) and not item.endswith('.txt'):
            files.append(item)
        elif os.path.isdir(item_path):
            folders.append(item)
    
    return render_template('home.html', files=files, folders=folders, current_path=path)

@app.errorhandler(403)
@app.errorhandler(404)

@app.route('/list_files')
def list_files():
    path = request.args.get('path', '')
    full_path = os.path.abspath(os.path.join(app.config['UPLOAD_FOLDER'], path))

    # Ensure the path is within the UPLOAD_FOLDER
    if not full_path.startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.exists(full_path) or not os.path.isdir(full_path):
        return jsonify({'error': 'Folder not found'}), 404

    files = []
    folders = []
    for item in os.listdir(full_path):
        item_path = os.path.join(full_path, item)
        if os.path.isfile(item_path) and not item.endswith('.txt'):
            file_details = {
                'name': item,
                'size': f"{os.path.getsize(item_path) / 1024:.2f} KB",
                'date_uploaded': datetime.fromtimestamp(os.path.getmtime(item_path)).strftime('%Y-%m-%d %H:%M:%S')
            }
            files.append(file_details)
        elif os.path.isdir(item_path):
            folders.append(item)

    return render_template('file_list.html', files=files, folders=folders, current_path=path)

def handle_error(error):
    return render_template('error.html', error=error.description), error.code

@app.route('/create_folder', methods=['POST'])
@app.route('/create_folder', methods=['POST'])
def create_folder():
    data = request.get_json()
    folder_name = data.get('folder_name')
    parent_path = data.get('parent_path', '')
    if folder_name:
        new_folder_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(parent_path), secure_filename(folder_name))
        try:
            os.makedirs(new_folder_path, exist_ok=True)
            return jsonify({'message': 'Folder created successfully'}), 200
        except Exception as e:
            return jsonify({'error': f'Failed to create folder: {str(e)}'}), 500
    return jsonify({'error': 'Invalid folder name'}), 400

@app.route('/paste', methods=['POST'])
def paste_items():
    data = request.get_json()
    clipboard = data.get('clipboard', {})
    destination_path = data.get('destination_path', '')
    
    if not clipboard:
        return jsonify({'error': 'Invalid clipboard data'}), 400
    
    action = clipboard['action']
    items = clipboard['items']
    source_path = clipboard['sourcePath']
    
    for item in items:
        source = os.path.join(app.config['UPLOAD_FOLDER'], source_path, item['name'])
        dest = os.path.join(app.config['UPLOAD_FOLDER'], destination_path, item['name'])
        
        if action == 'copy':
            if item['type'] == 'file':
                shutil.copy2(source, dest)
                # Copy associated .txt file if it exists
                txt_source = f"{source}.txt"
                txt_dest = f"{dest}.txt"
                if os.path.exists(txt_source):
                    shutil.copy2(txt_source, txt_dest)
            elif item['type'] == 'folder':
                shutil.copytree(source, dest)
        elif action == 'cut':
            shutil.move(source, dest)
            # Move associated .txt file if it exists
            txt_source = f"{source}.txt"
            txt_dest = f"{dest}.txt"
            if os.path.exists(txt_source):
                shutil.move(txt_source, txt_dest)
    
    return jsonify({'message': f'Items {action}d successfully'}), 200

@app.route('/delete', methods=['POST'])
def delete_items():
    data = request.get_json()
    items = data.get('items', [])
    parent_path = data.get('parent_path', '')
    
    for item in items:
        item_path = os.path.join(app.config['UPLOAD_FOLDER'], parent_path, item)
        if os.path.isfile(item_path):
            os.remove(item_path)
            # Remove associated .txt file if it exists
            txt_path = f"{item_path}.txt"
            if os.path.exists(txt_path):
                os.remove(txt_path)
        elif os.path.isdir(item_path):
            shutil.rmtree(item_path)
    
    return jsonify({'message': 'Items deleted successfully'}), 200

@app.route('/download/<path:filename>')
def download_file(filename):
    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)

@app.route('/download_converted/<path:filename>')
def download_converted(filename):
    converted_file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{filename}.txt")
    if os.path.exists(converted_file_path):
        return send_file(converted_file_path, as_attachment=True, download_name=f"{os.path.basename(filename)}_converted.txt")
    return jsonify({'error': 'Converted file not found'}), 404

@app.route('/upload', methods=['POST'])

def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    parent_path = request.form.get('parent_path', '')
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if file:
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], parent_path, filename)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        file.save(file_path)

        # Save file details
        file_details = {
            'name': file.filename,
            'type': file.filename.split('.')[-1],  # Get the file extension
            'size': f"{os.path.getsize(file_path) / 1024:.2f} KB",
            'date_uploaded': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'filepath': file_path
        }
        file_data_store[file.filename] = file_details
        return jsonify({'message': 'File uploaded successfully!', 'file_details': file_details})
        # Extract and store content for searching
        try:
            content = process_document(file_path)
            with open(f"{file_path}.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            return jsonify({'message': 'File uploaded successfully'}), 200
        except Exception as e:
            print(f"Error processing file: {e}")
            return jsonify({'error': 'Error processing file'}), 500
@app.route('/view_file/<path:filename>')
def view_file(filename):
    try:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Validate the file path is within UPLOAD_FOLDER
        if not os.path.abspath(file_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
            return jsonify({'error': 'Access denied'}), 403
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
        
        # Get file extension
        ext = os.path.splitext(filename)[1].lower()
        
        # Handle different file types
        if ext in ['.pdf', '.png', '.jpg', '.jpeg', '.gif']:
            return send_file(file_path)
            
        elif ext in ['.xlsx', '.xls']:
            # Convert Excel to HTML for viewing
            df = pd.read_excel(file_path)
            html_content = df.to_html(classes='table table-striped', index=False)
            return html_content
            
        elif ext in ['.docx', '.doc']:
            # Convert DOC/DOCX to HTML for viewing
            doc = Document(file_path)
            html_content = []
            for para in doc.paragraphs:
                html_content.append(f'<p>{para.text}</p>')
            return '\n'.join(html_content)
            
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
            
    except Exception as e:
        print(f"Error viewing file: {e}")
        return jsonify({'error': 'Error processing file'}), 500
@app.route('/search')
def search():
    query = request.args.get('query', '').lower()
    search_type = request.args.get('type', 'name')
    current_path = request.args.get('current_path', '')
    results = []

    search_root = os.path.join(app.config['UPLOAD_FOLDER'], current_path)
    
    for root, dirs, files in os.walk(search_root):
        for item in files:
            if item.endswith('.txt'):
                continue
            item_path = os.path.join(root, item)
            relative_path = os.path.relpath(item_path, search_root)
            
            if search_type == 'name':
                if query in item.lower():
                    results.append({'name': item, 'path': relative_path, 'type': 'file'})
            elif search_type == 'content':
                converted_file_path = f"{item_path}.txt"
                if os.path.exists(converted_file_path):
                    with open(converted_file_path, 'r', encoding='utf-8') as f:
                        content = f.read().lower()
                        if query in content:
                            results.append({'name': item, 'path': relative_path, 'type': 'file'})

    return jsonify(results)
@app.route('/extract_text/<path:filename>')
def extract_text(filename):
    try:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        print(f"Processing file: {file_path}")  # Debug log
        
        # Validate file path
        if not os.path.abspath(file_path).startswith(os.path.abspath(app.config['UPLOAD_FOLDER'])):
            return jsonify({'error': 'Access denied'}), 403
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404

        ext = os.path.splitext(filename)[1].lower()
        
        if ext == '.pdf':
            try:
                # Convert PDF to images
                images = convert_from_path(file_path, poppler_path=POPPLER_PATH)
                text_blocks = []
                
                for page_num, image in enumerate(images, 1):
                    try:
                        # Get OCR data with bounding boxes
                        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                        
                        for i in range(len(data['text'])):
                            if data['text'][i].strip():
                                text_blocks.append({
                                    'text': data['text'][i],
                                    'bbox': {
                                        'x': data['left'][i],
                                        'y': data['top'][i] + (page_num - 1) * image.size[1],
                                        'width': data['width'][i],
                                        'height': data['height'][i]
                                    },
                                    'page': page_num
                                })
                        
                    except Exception as e:
                        print(f"Error processing PDF page {page_num}: {e}")
                        continue
                
                return jsonify({'text': text_blocks})
                
            except Exception as e:
                print(f"Error processing PDF: {e}")
                return jsonify({'error': f'PDF processing error: {str(e)}'}), 500
            
        elif ext in ['.png', '.jpg', '.jpeg', '.gif']:
            try:
                # Process image
                with Image.open(file_path) as image:
                    # Convert image to RGB if necessary
                    if image.mode not in ('L', 'RGB'):
                        image = image.convert('RGB')
                    
                    # Get OCR data
                    data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                    
                    text_blocks = []
                    for i in range(len(data['text'])):
                        if data['text'][i].strip():
                            text_blocks.append({
                                'text': data['text'][i],
                                'bbox': {
                                    'x': data['left'][i],
                                    'y': data['top'][i],
                                    'width': data['width'][i],
                                    'height': data['height'][i]
                                }
                            })
                    
                    return jsonify({'text': text_blocks})
                    
            except Exception as e:
                print(f"Error processing image: {e}")
                return jsonify({'error': f'Image processing error: {str(e)}'}), 500
            
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
            
    except Exception as e:
        print(f"Error extracting text: {e}")
        print(traceback.format_exc())  # Print full traceback for debugging
        return jsonify({'error': str(e)}), 500
    

# Route to handle language selection
@app.route('/set_language', methods=['POST'])
def set_language():
    selected_language = request.form.get('language')
    return get_translation(selected_language)

@app.route('/reminder')
def reminder_page():
    return render_template('reminder.html')  # Reminder page with the form

@app.route('/reminder/add_reminder', methods=['POST'])
def add_reminder():
    data = request.get_json()
    title = data.get('title')
    date = data.get('date')
    time = data.get('time')
    toggle = data.get('toggle')

    # Call the reminder function here
    response = add_reminder(title, date, time, toggle)
    return jsonify(response)  # Return the response from the function

if __name__ == '__main__':
    app.run(debug=True)


